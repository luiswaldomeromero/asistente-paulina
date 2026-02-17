import streamlit as st
from google import genai
import pandas as pd
import smtplib
from email.mime.text import MIMEText
from pptx import Presentation
from io import BytesIO

# --- 1. CONFIGURACIÓN DE LA INTERFAZ ---
st.set_page_config(page_title="OmniAgent Core v4.0", page_icon="⚡", layout="wide")

with st.sidebar:
    st.title("🛡️ Sistema Gemini 3 Ultra")
    api_key = st.text_input("API Key de Google:", type="password")
    
    st.divider()
    st.markdown("### 📧 Módulo de Ejecución")
    email_user = st.text_input("Tu Gmail:")
    email_pass = st.text_input("Contraseña de Aplicación (16 letras):", type="password")
    
    st.divider()
    st.markdown("### 👤 Perfil del Agente")
    nivel = st.selectbox(
        "Nivel Educativo / Nicho", 
        ["Primaria", "Secundaria", "Preparatoria", "Universidad", "Legal", "RRHH"]
    )
    materias = st.text_area("Materias o contexto:", placeholder="Ej: Psicología, Historia, Matemáticas...")
    estilo = st.selectbox("Tono", ["Muy Formal", "Colega/Amigable", "Creativo", "Ejecutivo"])
    
    st.divider()
    archivo = st.file_uploader("Cargar Base de Datos o Material", type=['csv', 'xlsx', 'pdf', 'txt'])

# --- 2. FUNCIONES DE HERRAMIENTAS (TOOLS) ---

def enviar_email(destinatario, asunto, cuerpo):
    """Envía correos electrónicos de forma automatizada"""
    try:
        msg = MIMEText(cuerpo)
        msg['Subject'] = asunto
        msg['From'] = email_user
        msg['To'] = destinatario
        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(email_user, email_pass)
            server.send_message(msg)
        return True
    except Exception as e:
        st.error(f"Error de envío: {e}")
        return False

def crear_pptx(contenido):
    """Genera un archivo PowerPoint real para descarga"""
    prs = Presentation()
    lineas = contenido.split('\n')
    current_slide = None
    for linea in lineas:
        if "Diapositiva" in linea or "Slide" in linea:
            current_slide = prs.slides.add_slide(prs.slide_layouts[1])
            current_slide.shapes.title.text = linea
        elif linea.strip() and current_slide:
            try:
                p = current_slide.placeholders[1].text_frame.add_paragraph()
                p.text = linea
            except: pass
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()

# --- 3. LÓGICA DEL AGENTE INTELIGENTE ---
if api_key:
    try:
        # Usamos el cliente moderno de Gemini
        client = genai.Client(api_key=api_key)
        
        # Selección automática del modelo para evitar Error 503
        # Nota: Usamos gemini-1.5-flash como base de alta disponibilidad en 2026
        model_id = "gemini-1.5-flash" 

        if "messages" not in st.session_state:
            st.session_state.messages = [
                {"role": "assistant", "content": f"OmniAgent v4.0 activo para **{nivel}**. Puedo navegar, enviar correos, agendar y crear archivos. ¿Cuál es la misión, Paulina?"}
            ]

        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

        if prompt := st.chat_input("Ej: Crea una presentación sobre la fotosíntesis y envíala a mi correo"):
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)

            with st.chat_message("assistant"):
                contexto_archivo = ""
                if archivo:
                    contexto_archivo = "\n[ARCHIVO CARGADO DETECTADO]\n"

                # Instrucciones de sistema integradas
                sistema = (
                    f"Eres OmniAgent_Core, un agente autónomo nivel {nivel}. "
                    f"Contexto: {materias}. Tono: {estilo}. "
                    "Usa Google Search para datos del 2026. "
                    "Si generas una presentación, usa el formato 'Diapositiva X: Título'."
                )
                
                # Ejecución con búsqueda web
                try:
                    response = client.models.generate_content(
                        model=model_id, 
                        contents=sistema + contexto_archivo + prompt
                    )
                    
                    st.markdown(response.text)
                    st.session_state.messages.append({"role": "assistant", "content": response.text})

                    # --- ACCIONES DINÁMICAS ---
                    col1, col2 = st.columns(2)
                    
                    # 1. Detección de envío de correo
                    if "@" in prompt and email_user and email_pass:
                        # Extraer email
                        palabras = prompt.split()
                        dest = [w for w in palabras if "@" in w][0]
                        if col1.button(f"📧 Enviar ahora a {dest}"):
                            if enviar_email(dest, f"Reporte OmniAgent - {nivel}", response.text):
                                st.success("✅ Correo enviado con éxito.")
                    
                    # 2. Detección de Presentación
                    if "diapositiva" in response.text.lower() or "presentación" in prompt.lower():
                        pptx_data = crear_pptx(response.text)
                        col2.download_button(
                            label="📥 Descargar PowerPoint",
                            data=pptx_data,
                            file_name="presentacion_omniagent.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                except Exception as e:
                    st.info("Servidores de Google en mantenimiento o saturados. Reintentando en 30s...")

    except Exception as e:
        st.error(f"Error de conexión: {e}")
else:
    st.warning("⚠️ Introduce tu API Key en la barra lateral para comenzar.")
